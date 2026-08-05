from __future__ import annotations

import asyncio
import hashlib
import json
import io
import re
import time
from datetime import datetime, timezone
from pathlib import Path
from typing import Any, List

from runtime.document_drafts import (
    add_citation,
    add_section_block,
    create_draft_record,
    draft_stats,
    draft_to_markdown,
    inspect_draft_record,
    load_draft,
    save_draft,
)
from runtime.browser_runtime import playwright_chromium_readiness
from runtime.tool_registry import ToolRegistry, ToolSpec


def _resolve_output_path(input_data: dict[str, Any], context: Any, default_title: str, ext: str) -> "Path":
    """从多种可能字段名中灵活解析输出路径。

    未找到路径时，根据标题在工作区根目录自动生成。"""
    raw_path = (
        input_data.get("path")
        or input_data.get("output_path")
        or input_data.get("file_path")
        or input_data.get("filename")
    )
    if not raw_path:
        title = input_data.get("title", default_title) or default_title
        safe_name = re.sub(r'[<>:"/\\|?*]', '', title)[:60].strip() or default_title
        raw_path = f"{safe_name}{ext}"
    path = context.path_guard.resolve(raw_path)
    if path.suffix.lower() != ext:
        path = path.with_suffix(ext)
    return path


def _backup_output(path: Path, context: Any) -> None:
    backup_file = getattr(context, "backup_file", None)
    if callable(backup_file):
        backup_file(path)


async def extract_docx_outline(input_data: dict[str, Any], context: Any) -> dict[str, Any]:
    path = context.path_guard.resolve(input_data.get("path"))
    suffix = path.suffix.lower()
    if suffix not in (".docx", ".doc"):
        raise ValueError("仅支持 .docx 和 .doc 格式")

    from .docx_parser import DocxParser
    parser = DocxParser()

    extract_text = input_data.get("extract_text", True)
    result = await parser.parse(
        path,
        extract_text=extract_text,
        extract_headings=True,
    )

    output: dict[str, Any] = {
        "path": str(path),
        "paragraph_count": result.paragraph_count,
        "headings": result.headings,
        "strategy": result.strategy,
    }
    if extract_text:
        output["text"] = result.text
        output["text_chars"] = len(result.text or "")
    if result.table_count > 0:
        output["table_count"] = result.table_count
    if result.warnings:
        output["warnings"] = result.warnings
    return output


async def extract_pdf_text_preview(input_data: dict[str, Any], context: Any) -> dict[str, Any]:
    path = context.path_guard.resolve(input_data.get("path"))
    if path.suffix.lower() != ".pdf":
        raise ValueError("only .pdf files are supported")

    max_pages = int(input_data.get("max_pages", 50))

    from .pdf_parser import PDFParser, ParseResult
    parser = PDFParser()
    result: ParseResult = await parser.parse(path, max_pages=max_pages, context=context)

    output: dict[str, Any] = {
        "path": str(path),
        "page_count": result.total_pages,
        "pages_parsed": result.pages_parsed,
        "text": result.text,
        "strategy": result.strategy,
        "garbled_ratio": round(result.garbled_ratio, 4),
        "ocr_used": result.ocr_used,
    }
    if result.ocr_pages > 0:
        output["ocr_pages"] = result.ocr_pages
    if result.cid_garbled:
        output["cid_garbled_detected"] = True
    if result.warnings:
        output["warnings"] = result.warnings
    return output


def _resolve_pdf_docx_output_path(input_data: dict[str, Any], context: Any, pdf_path: Path) -> Path:
    raw_path = (
        input_data.get("output_path")
        or input_data.get("target_path")
        or input_data.get("docx_path")
    )
    if not raw_path:
        raw_path = str(pdf_path.with_name(f"{pdf_path.stem}_extracted.docx"))
    path = context.path_guard.resolve(raw_path)
    if path.suffix.lower() != ".docx":
        path = path.with_suffix(".docx")
    return path


def _normalize_pdf_docx_mode(input_data: dict[str, Any]) -> str:
    raw = str(
        input_data.get("mode")
        or input_data.get("conversion_mode")
        or input_data.get("layout_mode")
        or "text_only"
    ).strip().lower()
    aliases = {
        "text": "text_only",
        "text_only": "text_only",
        "plain_text": "text_only",
        "text_with_images": "text_with_images",
        "with_images": "text_with_images",
        "images": "text_with_images",
        "image_order": "text_with_images",
    }
    return aliases.get(raw, "text_only")


def _pdf_block_text(block: dict[str, Any]) -> str:
    lines: list[str] = []
    for line in block.get("lines") or []:
        spans = line.get("spans") if isinstance(line, dict) else None
        if not spans:
            continue
        text = "".join(str(span.get("text") or "") for span in spans if isinstance(span, dict)).strip()
        if text:
            lines.append(text)
    return "\n".join(lines).strip()


def _image_bytes_for_docx(raw: bytes) -> bytes:
    try:
        from PIL import Image
    except ImportError:
        return raw
    try:
        image = Image.open(io.BytesIO(raw))
        if image.format in {"PNG", "JPEG", "JPG"}:
            return raw
        out = io.BytesIO()
        image.convert("RGB").save(out, format="PNG")
        return out.getvalue()
    except Exception:
        return raw


def _pdf_ordered_page_blocks(page: Any) -> tuple[list[dict[str, Any]], int, int]:
    page_dict = page.get_text("dict")
    blocks: list[dict[str, Any]] = []
    text_block_count = 0
    image_count = 0
    for raw_block in page_dict.get("blocks") or []:
        if not isinstance(raw_block, dict):
            continue
        bbox = raw_block.get("bbox") or [0, 0, 0, 0]
        sort_key = (
            float(bbox[1]) if len(bbox) > 1 else 0.0,
            float(bbox[0]) if len(bbox) > 0 else 0.0,
        )
        block_type = raw_block.get("type")
        if block_type == 0:
            text = _pdf_block_text(raw_block)
            if text:
                blocks.append({"type": "text", "text": text, "bbox": bbox, "sort_key": sort_key})
                text_block_count += 1
        elif block_type == 1:
            image_bytes = raw_block.get("image")
            if isinstance(image_bytes, bytes) and image_bytes:
                blocks.append({
                    "type": "image",
                    "image": _image_bytes_for_docx(image_bytes),
                    "bbox": bbox,
                    "sort_key": sort_key,
                    "ext": raw_block.get("ext") or "",
                })
                image_count += 1
    blocks.sort(key=lambda item: item["sort_key"])
    return blocks, text_block_count, image_count


def _log_pdf_docx_progress(context: Any, message: str, data: dict[str, Any]) -> None:
    try:
        progress_data = dict(data)
        progress_data.setdefault("kind", "pdf_to_docx")
        context.log("info", message, progress_data)
    except Exception:
        pass


def _should_log_page_progress(page_number: int, pages_total: int) -> bool:
    if page_number <= 1 or page_number >= pages_total:
        return True
    if pages_total <= 30:
        return page_number % 2 == 0
    return page_number % 5 == 0


def _extract_pdf_ordered_blocks_sync(pdf_path: Path, max_pages: int) -> dict[str, Any]:
    try:
        import fitz
    except ImportError as exc:
        raise RuntimeError("PyMuPDF is required for PDF text_with_images conversion") from exc

    pages: list[dict[str, Any]] = []
    text_block_count = 0
    image_count = 0
    warnings: list[str] = []

    with fitz.open(str(pdf_path)) as pdf:
        total_pages = len(pdf)
        limit = min(total_pages, max_pages) if max_pages > 0 else total_pages
        for page_index in range(limit):
            page = pdf[page_index]
            blocks, page_text_blocks, page_images = _pdf_ordered_page_blocks(page)
            text_block_count += page_text_blocks
            image_count += page_images
            pages.append({
                "page_number": page_index + 1,
                "width": float(page.rect.width),
                "height": float(page.rect.height),
                "blocks": blocks,
            })

    return {
        "total_pages": total_pages,
        "pages_parsed": limit,
        "pages": pages,
        "text_block_count": text_block_count,
        "image_count": image_count,
        "warnings": warnings,
    }


async def _extract_pdf_ordered_blocks(pdf_path: Path, max_pages: int) -> dict[str, Any]:
    return await asyncio.to_thread(_extract_pdf_ordered_blocks_sync, pdf_path, max_pages)


def _extract_pdf_to_docx_text_with_images_sync(
    *,
    pdf_path: Path,
    output_path: Path,
    title: str,
    max_pages: int,
    context: Any,
) -> dict[str, Any]:
    try:
        import fitz
    except ImportError as exc:
        raise RuntimeError("PyMuPDF is required for PDF text_with_images conversion") from exc
    try:
        from docx import Document
        from docx.oxml.ns import qn
        from docx.shared import Inches, Pt
    except ImportError as exc:
        raise RuntimeError("python-docx is required for document.extract_pdf_to_docx") from exc

    doc = Document()
    normal = doc.styles["Normal"]
    normal.font.name = "Microsoft YaHei"
    normal._element.rPr.rFonts.set(qn("w:eastAsia"), "Microsoft YaHei")
    normal.font.size = Pt(10.5)

    doc.add_heading(title, level=1)
    meta = doc.add_paragraph()
    meta.add_run(f"源文件：{pdf_path.name}\n")

    max_image_width = 6.0
    skipped_images = 0
    text_block_count = 0
    image_count = 0
    total_pages = 0
    pages_parsed = 0
    warnings: list[str] = []

    with fitz.open(str(pdf_path)) as pdf:
        total_pages = len(pdf)
        limit = min(total_pages, max_pages) if max_pages > 0 else total_pages
        meta.add_run(
            f"页数：{total_pages}；已解析：{limit}；模式：text_with_images；说明：按页面块顺序近似保留文字和图片。"
        )
        _log_pdf_docx_progress(
            context,
            f"pdf conversion started 0/{limit}",
            {
                "phase": "started",
                "pages_done": 0,
                "pages_total": limit,
                "source_pages": total_pages,
                "mode": "text_with_images",
            },
        )

        for page_index in range(limit):
            page = pdf[page_index]
            page_number = page_index + 1
            page_width = max(float(page.rect.width), 1.0)
            blocks, page_text_blocks, page_images = _pdf_ordered_page_blocks(page)
            text_block_count += page_text_blocks
            image_count += page_images
            pages_parsed = page_number

            if page_index > 0:
                doc.add_page_break()
            doc.add_heading(f"第 {page_number} 页", level=2)
            for block in blocks:
                if block.get("type") == "text":
                    text = str(block.get("text") or "").strip()
                    if not text:
                        continue
                    for part in text.split("\n"):
                        if part.strip():
                            doc.add_paragraph(part.strip())
                    continue

                if block.get("type") == "image":
                    image_bytes = block.get("image")
                    if not isinstance(image_bytes, bytes) or not image_bytes:
                        continue
                    bbox = block.get("bbox") or [0, 0, page_width, 0]
                    raw_width = max(float(bbox[2]) - float(bbox[0]), 1.0) if len(bbox) >= 3 else page_width
                    width_inches = max(1.0, min(max_image_width, (raw_width / page_width) * max_image_width))
                    try:
                        doc.add_picture(io.BytesIO(image_bytes), width=Inches(width_inches))
                    except Exception as exc:
                        skipped_images += 1
                        context.log("warning", f"pdf image skipped: {str(exc)[:160]}", {"page": page_number})

            if _should_log_page_progress(page_number, limit):
                _log_pdf_docx_progress(
                    context,
                    f"pdf page converted {page_number}/{limit}",
                    {
                        "phase": "progress",
                        "pages_done": page_number,
                        "pages_total": limit,
                        "source_pages": total_pages,
                        "text_block_count": text_block_count,
                        "image_count": image_count,
                        "skipped_image_count": skipped_images,
                    },
                )

    _backup_output(output_path, context)
    _log_pdf_docx_progress(
        context,
        f"pdf docx saving {pages_parsed}/{pages_parsed or total_pages}",
        {
            "phase": "saving",
            "pages_done": pages_parsed,
            "pages_total": pages_parsed or total_pages,
            "source_pages": total_pages,
            "text_block_count": text_block_count,
            "image_count": image_count,
            "skipped_image_count": skipped_images,
        },
    )
    doc.save(str(output_path))
    _log_pdf_docx_progress(
        context,
        f"pdf docx saved {pages_parsed}/{pages_parsed or total_pages}",
        {
            "phase": "saved",
            "pages_done": pages_parsed,
            "pages_total": pages_parsed or total_pages,
            "source_pages": total_pages,
            "text_block_count": text_block_count,
            "image_count": image_count,
            "skipped_image_count": skipped_images,
            "file_size": output_path.stat().st_size if output_path.exists() else 0,
        },
    )

    if skipped_images:
        warnings.append(f"skipped_images:{skipped_images}")
    return {
        "path": str(output_path.resolve()),
        "source_path": str(pdf_path.resolve()),
        "mode": "text_with_images",
        "page_count": total_pages,
        "pages_parsed": pages_parsed,
        "strategy": "pymupdf_blocks",
        "text_block_count": text_block_count,
        "image_count": image_count,
        "skipped_image_count": skipped_images,
        "file_size": output_path.stat().st_size if output_path.exists() else 0,
        "warnings": warnings,
    }


async def _extract_pdf_to_docx_text_with_images(
    *,
    pdf_path: Path,
    output_path: Path,
    title: str,
    max_pages: int,
    context: Any,
) -> dict[str, Any]:
    return await asyncio.to_thread(
        _extract_pdf_to_docx_text_with_images_sync,
        pdf_path=pdf_path,
        output_path=output_path,
        title=title,
        max_pages=max_pages,
        context=context,
    )


async def extract_pdf_to_docx(input_data: dict[str, Any], context: Any) -> dict[str, Any]:
    pdf_path = context.path_guard.resolve(input_data.get("path"))
    if pdf_path.suffix.lower() != ".pdf":
        raise ValueError("only .pdf files are supported")

    output_path = _resolve_pdf_docx_output_path(input_data, context, pdf_path)
    max_pages = int(input_data.get("max_pages", 0))
    title = str(input_data.get("title") or f"{pdf_path.stem} 提取文本").strip()
    mode = _normalize_pdf_docx_mode(input_data)
    if mode == "text_with_images":
        return await _extract_pdf_to_docx_text_with_images(
            pdf_path=pdf_path,
            output_path=output_path,
            title=title,
            max_pages=max_pages,
            context=context,
        )

    from .pdf_parser import PDFParser, ParseResult
    parser = PDFParser()
    result: ParseResult = await parser.parse(pdf_path, max_pages=max_pages, context=context)

    try:
        from docx import Document
        from docx.oxml.ns import qn
        from docx.shared import Pt
    except ImportError as exc:
        raise RuntimeError("python-docx is required for document.extract_pdf_to_docx") from exc

    doc = Document()
    normal = doc.styles["Normal"]
    normal.font.name = "Microsoft YaHei"
    normal._element.rPr.rFonts.set(qn("w:eastAsia"), "Microsoft YaHei")
    normal.font.size = Pt(10.5)

    doc.add_heading(title, level=1)
    meta = doc.add_paragraph()
    meta.add_run(f"源文件：{pdf_path.name}\n")
    meta.add_run(f"页数：{result.total_pages}；已解析：{result.pages_parsed}；策略：{result.strategy}")
    if result.ocr_used:
        meta.add_run(f"；OCR 页数：{result.ocr_pages}")

    text = (result.text or "").replace("\r\n", "\n").replace("\r", "\n").strip()
    if not text:
        doc.add_paragraph("未提取到可写入的文本内容。")
    else:
        for block in text.split("\n\n"):
            block = block.strip()
            if not block:
                continue
            for start in range(0, len(block), 3000):
                doc.add_paragraph(block[start:start + 3000])

    _backup_output(output_path, context)
    doc.save(str(output_path))

    output: dict[str, Any] = {
        "path": str(output_path.resolve()),
        "source_path": str(pdf_path.resolve()),
        "page_count": result.total_pages,
        "pages_parsed": result.pages_parsed,
        "strategy": result.strategy,
        "garbled_ratio": round(result.garbled_ratio, 4),
        "ocr_used": result.ocr_used,
    }
    if result.ocr_pages > 0:
        output["ocr_pages"] = result.ocr_pages
    if result.warnings:
        output["warnings"] = result.warnings
    return output


async def export_markdown(input_data: dict[str, Any], context: Any) -> dict[str, Any]:
    title = input_data.get("title", "AI生成文档")
    path = _resolve_output_path(input_data, context, title, ".md")
    content = input_data.get("content", "")

    content = content.strip()
    if not content.startswith("#"):
        content = f"# {title}\n\n{content}"

    _backup_output(path, context)
    path.write_text(content + "\n", encoding="utf-8")

    return {
        "path": str(path.resolve()),
        "size": len(content),
        "title": title,
    }


async def export_docx(input_data: dict[str, Any], context: Any) -> dict[str, Any]:
    title = input_data.get("title", "AI生成文档")
    path = _resolve_output_path(input_data, context, title, ".docx")
    content = input_data.get("content", "")

    try:
        from docx import Document
        from docx.oxml.ns import qn
        from docx.shared import Pt, RGBColor
    except ImportError as exc:
        raise RuntimeError("python-docx is required for document.export_docx") from exc

    doc = Document()

    normal = doc.styles["Normal"]
    normal.font.name = "Microsoft YaHei"
    normal._element.rPr.rFonts.set(qn("w:eastAsia"), "Microsoft YaHei")
    normal.font.size = Pt(10.5)

    for style_name in ("Heading 1", "Heading 2", "Heading 3"):
        style = doc.styles[style_name]
        style.font.name = "Microsoft YaHei"
        style._element.rPr.rFonts.set(qn("w:eastAsia"), "Microsoft YaHei")

    markdown = content.strip()
    if not markdown.lstrip().startswith("#"):
        doc.add_heading(title, level=1)

    lines = markdown.replace("\r\n", "\n").replace("\r", "\n").split("\n")
    index = 0
    in_code = False
    code_lines: List[str] = []

    while index < len(lines):
        line = lines[index].rstrip()
        stripped = line.strip()

        if stripped.startswith("```"):
            if in_code:
                p = doc.add_paragraph()
                p.style = "No Spacing"
                run = p.add_run("\n".join(code_lines))
                font = run.font
                font.name = "Consolas"
                code_lines = []
                in_code = False
            else:
                in_code = True
            index += 1
            continue

        if in_code:
            code_lines.append(line)
            index += 1
            continue

        if not stripped:
            doc.add_paragraph()
            index += 1
            continue

        heading_match = re.match(r"^(#{1,6})\s+(.+)$", stripped)
        if heading_match:
            level = min(len(heading_match.group(1)), 3)
            doc.add_heading(heading_match.group(2), level=level)
            index += 1
            continue

        doc.add_paragraph(stripped)
        index += 1

    if in_code and code_lines:
        p = doc.add_paragraph()
        p.style = "No Spacing"
        run = p.add_run("\n".join(code_lines))
        font = run.font
        font.name = "Consolas"

    paragraph_count = len(doc.paragraphs)
    nonempty_paragraph_count = sum(1 for paragraph in doc.paragraphs if paragraph.text.strip())

    _backup_output(path, context)
    doc.save(str(path))

    return {
        "path": str(path.resolve()),
        "title": title,
        "content_chars": len(content.strip()),
        "paragraph_count": paragraph_count,
        "nonempty_paragraph_count": nonempty_paragraph_count,
        "file_size": path.stat().st_size if path.exists() else 0,
    }


def _document_draft_data_dir(context: Any) -> Path:
    settings = getattr(context, "settings", None)
    data_dir = getattr(settings, "data_dir", None)
    if data_dir:
        return Path(data_dir)
    temp_dir = getattr(context, "temp_dir", None)
    if temp_dir:
        return Path(temp_dir) / "runtime-data"
    raise RuntimeError("document draft tools require settings.data_dir or context.temp_dir")


def _safe_document_draft_level(value: Any) -> int:
    try:
        return int(value)
    except (TypeError, ValueError):
        return 1


def _draft_record_from_docx(source_path: Path, title: str = "") -> dict[str, Any]:
    try:
        from docx import Document
    except ImportError as exc:
        raise RuntimeError("python-docx is required to import a Word document into a draft") from exc

    doc = Document(str(source_path))
    document_title = str(title or doc.core_properties.title or source_path.stem).strip() or source_path.stem
    section_specs: list[dict[str, Any]] = []
    section_blocks: list[list[str]] = []
    current_index = -1

    for paragraph in doc.paragraphs:
        text = str(paragraph.text or "").strip()
        if not text:
            continue
        style_name = str(getattr(paragraph.style, "name", "") or "")
        heading_match = re.search(r"(?:heading|标题)\s*([1-6])", style_name, re.IGNORECASE)
        if heading_match:
            section_specs.append({
                "title": text,
                "level": int(heading_match.group(1)),
            })
            section_blocks.append([])
            current_index = len(section_specs) - 1
            continue
        if current_index < 0:
            section_specs.append({"section_id": "body", "title": "Body", "level": 1})
            section_blocks.append([])
            current_index = 0
        section_blocks[current_index].append(text)

    record = create_draft_record(
        title=document_title,
        sections=section_specs,
        metadata={
            "source_path": str(source_path.resolve()),
            "source_type": "docx",
            "imported": True,
        },
    )
    for section, blocks in zip(record.get("sections") or [], section_blocks):
        for text in blocks:
            add_section_block(
                record,
                section_id=str(section.get("section_id") or ""),
                content=text,
                metadata={"imported": True},
            )
    return record


async def create_draft(input_data: dict[str, Any], context: Any) -> dict[str, Any]:
    source_path_value = input_data.get("source_path")
    if source_path_value:
        source_path = context.path_guard.resolve(str(source_path_value))
        if source_path.suffix.lower() != ".docx":
            raise ValueError("document.create_draft source_path currently supports .docx files only")
        record = _draft_record_from_docx(source_path, str(input_data.get("title") or ""))
        record["metadata"].update(
            input_data.get("metadata") if isinstance(input_data.get("metadata"), dict) else {}
        )
    else:
        record = create_draft_record(
            title=str(input_data.get("title") or "Untitled Draft"),
            sections=input_data.get("sections") if isinstance(input_data.get("sections"), list) else [],
            metadata=input_data.get("metadata") if isinstance(input_data.get("metadata"), dict) else {},
        )
    data_dir = _document_draft_data_dir(context)
    record = save_draft(data_dir, record)
    return {
        "draft_id": record["draft_id"],
        "title": record["title"],
        "stats": draft_stats(record),
        "source_path": record.get("metadata", {}).get("source_path"),
        "message": (
            "draft created; append content in complete, bounded blocks with "
            "document.append_draft_section, inspect progress, then export with document.export_draft_docx"
        ),
    }


async def append_draft_section(input_data: dict[str, Any], context: Any) -> dict[str, Any]:
    data_dir = _document_draft_data_dir(context)
    record = load_draft(data_dir, str(input_data.get("draft_id") or ""))
    citation_ids = input_data.get("citation_ids")
    if not isinstance(citation_ids, list):
        citation_ids = []
    block = add_section_block(
        record,
        content=str(input_data.get("content") or ""),
        section_id=str(input_data.get("section_id") or "").strip() or None,
        title=str(input_data.get("title") or "").strip() or None,
        level=_safe_document_draft_level(input_data.get("level")),
        citation_ids=[str(item) for item in citation_ids],
        metadata=input_data.get("metadata") if isinstance(input_data.get("metadata"), dict) else {},
    )
    record = save_draft(data_dir, record)
    stats = draft_stats(record)
    return {
        "draft_id": record["draft_id"],
        "block_id": block["block_id"],
        "stats": stats,
        "unknown_citation_ids": stats.get("unknown_citation_ids", []),
    }


async def add_draft_citation(input_data: dict[str, Any], context: Any) -> dict[str, Any]:
    data_dir = _document_draft_data_dir(context)
    record = load_draft(data_dir, str(input_data.get("draft_id") or ""))
    citation = input_data.get("citation")
    if not isinstance(citation, dict):
        citation = {key: value for key, value in input_data.items() if key != "draft_id"}
    item = add_citation(record, citation)
    record = save_draft(data_dir, record)
    return {
        "draft_id": record["draft_id"],
        "citation": item,
        "stats": draft_stats(record),
    }


async def inspect_draft(input_data: dict[str, Any], context: Any) -> dict[str, Any]:
    data_dir = _document_draft_data_dir(context)
    record = load_draft(data_dir, str(input_data.get("draft_id") or ""))
    return inspect_draft_record(record)


async def export_draft_docx(input_data: dict[str, Any], context: Any) -> dict[str, Any]:
    data_dir = _document_draft_data_dir(context)
    record = load_draft(data_dir, str(input_data.get("draft_id") or ""))
    include_citations = bool(input_data.get("include_citations", True))
    markdown = draft_to_markdown(record, include_citations=include_citations)
    output_path = input_data.get("output_path") or input_data.get("path")
    output = await export_docx(
        {
            "path": output_path,
            "title": record.get("title") or "Document Draft",
            "content": markdown,
        },
        context,
    )
    output["draft_id"] = record["draft_id"]
    output["draft_stats"] = draft_stats(record)
    return output


def _resolve_translate_docx_output_path(input_data: dict[str, Any], context: Any, source_path: Path) -> Path:
    raw_path = (
        input_data.get("output_path")
        or input_data.get("target_path")
        or input_data.get("translated_path")
    )
    if not raw_path:
        raw_path = str(source_path.with_name(f"{source_path.stem}_translated_zh.docx"))
    path = context.path_guard.resolve(raw_path)
    if path.suffix.lower() != ".docx":
        path = path.with_suffix(".docx")
    return path


def _split_translation_chunks(text: str, max_chars: int) -> list[str]:
    if len(text) <= max_chars:
        return [text]
    chunks: list[str] = []
    current = ""
    parts = re.split(r"(?<=[.!?。！？；;])\s+", text)
    for part in parts:
        if not part:
            continue
        if len(part) > max_chars:
            if current:
                chunks.append(current)
                current = ""
            chunks.extend(part[index:index + max_chars] for index in range(0, len(part), max_chars))
            continue
        if current and len(current) + len(part) + 1 > max_chars:
            chunks.append(current)
            current = part
        else:
            current = f"{current} {part}".strip() if current else part
    if current:
        chunks.append(current)
    return chunks


async def translate_docx(input_data: dict[str, Any], context: Any) -> dict[str, Any]:
    engine = str(input_data.get("engine") or "model").strip().lower()
    if engine in {"model", "llm", "ai"}:
        return await _translate_docx_model(input_data, context)
    if engine in {"google", "google_translate", "external"}:
        return await asyncio.to_thread(_translate_docx_google_sync, input_data, context)
    raise ValueError("document.translate_docx engine must be 'model' or 'google'")


def _model_translation_prompt(source_language: str, target_language: str, texts: list[str]) -> list[dict[str, Any]]:
    items = [{"i": index, "text": text} for index, text in enumerate(texts)]
    payload = {
        "source_language": source_language,
        "target_language": target_language,
        "items": items,
    }
    return [
        {
            "role": "system",
            "content": (
                "You are a document translation engine. Translate every input item into the target language. "
                "Do not summarize, omit, add commentary, or change the item count. Preserve names, numbers, "
                "format-like markers, and paragraph meaning. Return only valid JSON."
            ),
        },
        {
            "role": "user",
            "content": (
                "Return a JSON array like "
                '[{"i":0,"text":"translated text"}]. '
                "Use the same i values from the input.\n"
                + json.dumps(payload, ensure_ascii=False)
            ),
        },
    ]


def _extract_json_payload(text: str) -> Any:
    value = text.strip()
    if value.startswith("```"):
        value = re.sub(r"^```(?:json)?\s*", "", value, flags=re.IGNORECASE)
        value = re.sub(r"\s*```$", "", value).strip()
    try:
        return json.loads(value)
    except json.JSONDecodeError:
        pass

    array_start = value.find("[")
    array_end = value.rfind("]")
    if array_start >= 0 and array_end > array_start:
        return json.loads(value[array_start:array_end + 1])

    object_start = value.find("{")
    object_end = value.rfind("}")
    if object_start >= 0 and object_end > object_start:
        return json.loads(value[object_start:object_end + 1])

    raise ValueError("model did not return valid JSON")


def _parse_translation_response(answer: str, expected_count: int) -> list[str]:
    data = _extract_json_payload(answer)
    if isinstance(data, dict):
        for key in ("translations", "items", "results", "result"):
            if isinstance(data.get(key), list):
                data = data[key]
                break

    if isinstance(data, list) and all(isinstance(item, str) for item in data):
        translations = [str(item).strip() for item in data]
    elif isinstance(data, list):
        by_index: dict[int, str] = {}
        ordered: list[str] = []
        for fallback_index, item in enumerate(data):
            if not isinstance(item, dict):
                raise ValueError("translation JSON items must be objects or strings")
            raw_index = item.get("i", item.get("index", item.get("id", fallback_index)))
            try:
                index = int(raw_index)
            except (TypeError, ValueError):
                index = fallback_index
            text = item.get("text", item.get("translation", item.get("translated_text")))
            if text is None:
                raise ValueError("translation JSON item is missing text")
            by_index[index] = str(text).strip()
            ordered.append(str(text).strip())
        translations = [by_index.get(index, ordered[index] if index < len(ordered) else "") for index in range(expected_count)]
    else:
        raise ValueError("translation JSON must be an array")

    if len(translations) != expected_count:
        raise ValueError(f"translation count mismatch: expected {expected_count}, got {len(translations)}")
    if any(not item for item in translations):
        raise ValueError("translation response contained empty text")
    return translations


async def _translate_text_batch_with_model(
    *,
    settings: Any,
    model: str,
    texts: list[str],
    source_language: str,
    target_language: str,
) -> list[str]:
    from runtime.model_providers.client import generate_chat_completion

    answer, _metadata = await generate_chat_completion(
        settings=settings,
        model=model,
        messages=_model_translation_prompt(source_language, target_language, texts),
        enable_thinking=False,
        reasoning_effort="low",
        tools=None,
    )
    return _parse_translation_response(answer, len(texts))


_TRANSLATE_DOCX_DEFAULT_PROFILE = "balanced"
_TRANSLATE_DOCX_PROFILES: dict[str, dict[str, int]] = {
    # 默认保持响应及时；大上下文模型可显式使用 "fast"。
    "safe": {
        "max_chars_per_chunk": 3000,
        "max_chars_per_batch": 6000,
        "max_paragraphs_per_batch": 4,
        "batch_timeout": 180,
    },
    "balanced": {
        "max_chars_per_chunk": 4000,
        "max_chars_per_batch": 12000,
        "max_paragraphs_per_batch": 8,
        "batch_timeout": 240,
    },
    "fast": {
        "max_chars_per_chunk": 6000,
        "max_chars_per_batch": 24000,
        "max_paragraphs_per_batch": 16,
        "batch_timeout": 300,
    },
}


def _resolve_translate_docx_profile(input_data: dict[str, Any]) -> tuple[str, dict[str, int]]:
    requested = str(
        input_data.get("translation_profile")
        or input_data.get("profile")
        or _TRANSLATE_DOCX_DEFAULT_PROFILE
    ).strip().lower()
    aliases = {
        "default": "balanced",
        "normal": "balanced",
        "stable": "safe",
        "large": "fast",
    }
    profile = aliases.get(requested, requested)
    if profile not in _TRANSLATE_DOCX_PROFILES:
        profile = _TRANSLATE_DOCX_DEFAULT_PROFILE
    return profile, _TRANSLATE_DOCX_PROFILES[profile]


def _translate_docx_int(input_data: dict[str, Any], key: str, default: int) -> int:
    try:
        return int(input_data.get(key, default) or default)
    except (TypeError, ValueError):
        return default


def _translate_docx_bool(input_data: dict[str, Any], key: str, default: bool) -> bool:
    value = input_data.get(key)
    if value is None:
        return default
    if isinstance(value, bool):
        return value
    if isinstance(value, (int, float)):
        return bool(value)
    text = str(value).strip().lower()
    if text in {"1", "true", "yes", "y", "on"}:
        return True
    if text in {"0", "false", "no", "n", "off"}:
        return False
    return default


_TRANSLATE_DOCX_CHECKPOINT_SCHEMA_VERSION = "0.1"
_TRANSLATE_DOCX_CHECKPOINT_KIND = "document.translate_docx.checkpoint"
_TRANSLATE_DOCX_DONE_STATUSES = {"translated", "failed"}


def _translate_docx_now() -> str:
    return datetime.now(timezone.utc).isoformat(timespec="seconds")


def _translate_docx_source_hash(text: str) -> str:
    return hashlib.sha256(text.encode("utf-8")).hexdigest()


def _resolve_translate_docx_manifest_path(input_data: dict[str, Any], context: Any, output_path: Path) -> Path:
    raw_path = input_data.get("manifest_path") or input_data.get("checkpoint_path")
    if raw_path:
        path = context.path_guard.resolve(raw_path)
    else:
        path = output_path.with_suffix(output_path.suffix + ".translate.json")
    if path.suffix.lower() != ".json":
        path = path.with_suffix(".json")
    return path


def _build_translate_docx_source_items(
    paragraphs: list[Any],
    target_nonempty_goal: int,
) -> list[dict[str, Any]]:
    items: list[dict[str, Any]] = []
    nonempty_index = 0
    for paragraph_index, paragraph in enumerate(paragraphs, 1):
        text = paragraph.text.strip()
        if not text:
            continue
        nonempty_index += 1
        if len(items) >= target_nonempty_goal:
            break
        items.append({
            "paragraph_index": paragraph_index,
            "nonempty_index": nonempty_index,
            "source_text": text,
            "source_hash": _translate_docx_source_hash(text),
            "source_chars": len(text),
        })
    return items


def _new_translate_docx_manifest(
    *,
    source_path: Path,
    output_path: Path,
    model: str,
    source_language: str,
    target_language: str,
    source_nonempty: int,
    target_nonempty_goal: int,
    source_items: list[dict[str, Any]],
) -> dict[str, Any]:
    now = _translate_docx_now()
    return {
        "schema_version": _TRANSLATE_DOCX_CHECKPOINT_SCHEMA_VERSION,
        "kind": _TRANSLATE_DOCX_CHECKPOINT_KIND,
        "source_path": str(source_path.resolve()),
        "output_path": str(output_path.resolve()),
        "model": model,
        "source_language": source_language,
        "target_language": target_language,
        "source_nonempty": source_nonempty,
        "target_nonempty_goal": target_nonempty_goal,
        "source_chars_total": sum(int(item.get("source_chars") or 0) for item in source_items),
        "created_at": now,
        "updated_at": now,
        "items": [
            {
                "paragraph_index": item["paragraph_index"],
                "nonempty_index": item["nonempty_index"],
                "source_hash": item["source_hash"],
                "source_chars": item["source_chars"],
                "source_text": item["source_text"],
                "status": "pending",
                "translated_text": "",
                "attempts": 0,
                "error": "",
                "updated_at": "",
            }
            for item in source_items
        ],
        "warnings": [],
    }


def _translate_docx_manifest_matches(
    manifest: dict[str, Any],
    *,
    source_items: list[dict[str, Any]],
    source_language: str,
    target_language: str,
    target_nonempty_goal: int,
) -> bool:
    if manifest.get("kind") != _TRANSLATE_DOCX_CHECKPOINT_KIND:
        return False
    if str(manifest.get("source_language") or "") != source_language:
        return False
    if str(manifest.get("target_language") or "") != target_language:
        return False
    if int(manifest.get("target_nonempty_goal") or 0) != target_nonempty_goal:
        return False
    items = manifest.get("items")
    if not isinstance(items, list) or len(items) != len(source_items):
        return False
    for saved, current in zip(items, source_items):
        if not isinstance(saved, dict):
            return False
        if int(saved.get("paragraph_index") or 0) != int(current.get("paragraph_index") or 0):
            return False
        if str(saved.get("source_hash") or "") != str(current.get("source_hash") or ""):
            return False
    return True


def _save_translate_docx_manifest(path: Path, manifest: dict[str, Any]) -> None:
    manifest["updated_at"] = _translate_docx_now()
    path.parent.mkdir(parents=True, exist_ok=True)
    path.write_text(json.dumps(manifest, ensure_ascii=False, indent=2), encoding="utf-8")


def _load_translate_docx_manifest(path: Path) -> dict[str, Any] | None:
    if not path.exists():
        return None
    try:
        value = json.loads(path.read_text(encoding="utf-8"))
    except (OSError, json.JSONDecodeError):
        return None
    return value if isinstance(value, dict) else None


def _translate_docx_completed_prefix_count(items: list[dict[str, Any]]) -> int:
    count = 0
    for item in items:
        if str(item.get("status") or "") not in _TRANSLATE_DOCX_DONE_STATUSES:
            break
        count += 1
    return count


def _append_source_blanks_until(
    target_doc: Any,
    paragraphs: list[Any],
    last_written_paragraph_index: int,
    next_paragraph_index: int,
) -> None:
    for paragraph_index in range(last_written_paragraph_index + 1, next_paragraph_index):
        paragraph = paragraphs[paragraph_index - 1]
        if not paragraph.text.strip():
            target_doc.add_paragraph()


def _append_translate_docx_completed_prefix(
    *,
    target_doc: Any,
    paragraphs: list[Any],
    manifest_items: list[dict[str, Any]],
) -> dict[str, int]:
    last_written_paragraph_index = 0
    translated_count = 0
    failed_count = 0
    translated_chars = 0
    completed_source_chars = 0
    prefix_count = _translate_docx_completed_prefix_count(manifest_items)
    for item in manifest_items[:prefix_count]:
        paragraph_index = int(item.get("paragraph_index") or 0)
        _append_source_blanks_until(target_doc, paragraphs, last_written_paragraph_index, paragraph_index)
        status = str(item.get("status") or "")
        if status == "translated":
            translated = str(item.get("translated_text") or "").strip()
            target_doc.add_paragraph(translated)
            translated_count += 1
            translated_chars += len(translated)
        elif status == "failed":
            original = str(item.get("source_text") or "")
            target_doc.add_paragraph(f"[translation failed, original kept] {original}")
            failed_count += 1
        completed_source_chars += int(item.get("source_chars") or 0)
        last_written_paragraph_index = paragraph_index
    return {
        "prefix_count": prefix_count,
        "last_written_paragraph_index": last_written_paragraph_index,
        "translated_count": translated_count,
        "failed_count": failed_count,
        "translated_chars": translated_chars,
        "completed_source_chars": completed_source_chars,
    }


def _translate_docx_load_or_create_manifest(
    *,
    input_data: dict[str, Any],
    context: Any,
    manifest_path: Path,
    source_path: Path,
    output_path: Path,
    model: str,
    source_language: str,
    target_language: str,
    source_nonempty: int,
    target_nonempty_goal: int,
    source_items: list[dict[str, Any]],
) -> tuple[dict[str, Any], bool]:
    resume = _translate_docx_bool(input_data, "resume", True)
    if _translate_docx_bool(input_data, "reset_checkpoint", False):
        resume = False
    manifest = _load_translate_docx_manifest(manifest_path) if resume else None
    if manifest and _translate_docx_manifest_matches(
        manifest,
        source_items=source_items,
        source_language=source_language,
        target_language=target_language,
        target_nonempty_goal=target_nonempty_goal,
    ):
        return manifest, True
    manifest = _new_translate_docx_manifest(
        source_path=source_path,
        output_path=output_path,
        model=model,
        source_language=source_language,
        target_language=target_language,
        source_nonempty=source_nonempty,
        target_nonempty_goal=target_nonempty_goal,
        source_items=source_items,
    )
    _save_translate_docx_manifest(manifest_path, manifest)
    return manifest, False


async def _translate_docx_model(input_data: dict[str, Any], context: Any) -> dict[str, Any]:
    source_path = context.path_guard.resolve(
        input_data.get("path") or input_data.get("source_path")
    )
    if source_path.suffix.lower() != ".docx":
        raise ValueError("only .docx files are supported")

    settings = getattr(context, "settings", None)
    model = str(input_data.get("model") or "").strip()
    if not model and settings and hasattr(settings, "get_default_model"):
        model = str(settings.get_default_model())
    if not settings or not model:
        raise RuntimeError("document.translate_docx model engine requires runtime settings and a model")

    output_path = _resolve_translate_docx_output_path(input_data, context, source_path)
    source_language = str(input_data.get("source_language") or input_data.get("source") or "auto")
    target_language = str(input_data.get("target_language") or input_data.get("target") or "zh-CN")
    translation_profile, profile_defaults = _resolve_translate_docx_profile(input_data)
    max_chars_per_chunk = max(
        1000,
        min(
            _translate_docx_int(input_data, "max_chars_per_chunk", profile_defaults["max_chars_per_chunk"]),
            12000,
        ),
    )
    max_chars_per_batch = max(
        max_chars_per_chunk,
        min(
            _translate_docx_int(input_data, "max_chars_per_batch", profile_defaults["max_chars_per_batch"]),
            64000,
        ),
    )
    max_paragraphs_per_batch = max(
        1,
        min(
            _translate_docx_int(
                input_data,
                "max_paragraphs_per_batch",
                profile_defaults["max_paragraphs_per_batch"],
            ),
            40,
        ),
    )
    batch_timeout = max(
        60,
        min(_translate_docx_int(input_data, "batch_timeout", profile_defaults["batch_timeout"]), 600),
    )
    max_paragraphs = max(0, int(input_data.get("max_paragraphs", 0) or 0))
    max_seconds = max(120, min(int(input_data.get("max_seconds", 1800) or 1800), 7200))
    if not max_paragraphs:
        max_seconds = max(max_seconds, 1800)
    save_every = max(1, min(int(input_data.get("save_every", 10) or 10), 100))

    try:
        from docx import Document
        from docx.oxml.ns import qn
        from docx.shared import Pt
    except ImportError as exc:
        raise RuntimeError("python-docx is required for document.translate_docx") from exc

    source_doc = Document(str(source_path))
    paragraphs = list(source_doc.paragraphs)
    source_texts = [paragraph.text.strip() for paragraph in paragraphs if paragraph.text.strip()]
    source_nonempty = len(source_texts)
    target_nonempty_goal = min(source_nonempty, max_paragraphs) if max_paragraphs else source_nonempty
    source_items = _build_translate_docx_source_items(paragraphs, target_nonempty_goal)
    target_source_chars = sum(int(item.get("source_chars") or 0) for item in source_items)
    manifest_path = _resolve_translate_docx_manifest_path(input_data, context, output_path)
    manifest, resumed_from_checkpoint = _translate_docx_load_or_create_manifest(
        input_data=input_data,
        context=context,
        manifest_path=manifest_path,
        source_path=source_path,
        output_path=output_path,
        model=model,
        source_language=source_language,
        target_language=target_language,
        source_nonempty=source_nonempty,
        target_nonempty_goal=target_nonempty_goal,
        source_items=source_items,
    )
    manifest_items = manifest.get("items") if isinstance(manifest.get("items"), list) else []

    target_doc = Document()
    normal = target_doc.styles["Normal"]
    normal.font.name = "Microsoft YaHei"
    normal._element.rPr.rFonts.set(qn("w:eastAsia"), "Microsoft YaHei")
    normal.font.size = Pt(10.5)

    prefix_state = _append_translate_docx_completed_prefix(
        target_doc=target_doc,
        paragraphs=paragraphs,
        manifest_items=manifest_items,
    )
    started_at = time.monotonic()
    processed_nonempty = int(prefix_state["prefix_count"])
    translated_count = int(prefix_state["translated_count"])
    failed_count = int(prefix_state["failed_count"])
    completed_source_chars = int(prefix_state["completed_source_chars"])
    translated_chars = int(prefix_state["translated_chars"])
    last_written_paragraph_index = int(prefix_state["last_written_paragraph_index"])
    warnings: list[str] = [
        str(item)
        for item in manifest.get("warnings", [])
        if isinstance(item, str)
    ]
    stopped_reason = ""
    pending: list[dict[str, Any]] = []
    pending_chars = 0

    _backup_output(output_path, context)
    context.log(
        "info",
        f"translation source loaded {processed_nonempty}/{target_nonempty_goal}",
        {
            "source_nonempty": source_nonempty,
            "target_goal": target_nonempty_goal,
            "source_chars_total": target_source_chars,
            "source_chars_done": completed_source_chars,
            "engine": "model",
            "translation_profile": translation_profile,
            "manifest_path": str(manifest_path.resolve()),
            "resumed_from_checkpoint": resumed_from_checkpoint,
            "max_paragraphs_per_batch": max_paragraphs_per_batch,
            "max_chars_per_batch": max_chars_per_batch,
            "max_chars_per_chunk": max_chars_per_chunk,
            "batch_timeout": batch_timeout,
        },
    )

    def update_manifest_warning(message: str) -> None:
        warnings.append(message)
        manifest["warnings"] = warnings[-100:]

    def append_failed_paragraph(item: dict[str, Any], exc: Exception) -> None:
        nonlocal failed_count, last_written_paragraph_index
        paragraph_index = int(item.get("paragraph_index") or 0)
        original_text = str(item.get("source_text") or "")
        _append_source_blanks_until(target_doc, paragraphs, last_written_paragraph_index, paragraph_index)
        failed_count += 1
        target_doc.add_paragraph(f"[translation failed, original kept] {original_text}")
        last_written_paragraph_index = paragraph_index
        item["status"] = "failed"
        item["attempts"] = int(item.get("attempts") or 0) + 1
        item["error"] = str(exc)[:500]
        item["updated_at"] = _translate_docx_now()
        update_manifest_warning(f"paragraph {paragraph_index}: {str(exc)[:200]}")

    def append_translated_paragraph(item: dict[str, Any], translated: str) -> None:
        nonlocal translated_count, translated_chars, last_written_paragraph_index
        paragraph_index = int(item.get("paragraph_index") or 0)
        _append_source_blanks_until(target_doc, paragraphs, last_written_paragraph_index, paragraph_index)
        target_doc.add_paragraph(translated)
        translated_count += 1
        translated_chars += len(translated)
        last_written_paragraph_index = paragraph_index
        item["status"] = "translated"
        item["translated_text"] = translated
        item["attempts"] = int(item.get("attempts") or 0) + 1
        item["error"] = ""
        item["updated_at"] = _translate_docx_now()

    def log_translation_progress(*, fallback: bool = False, paragraph_index: int | None = None) -> None:
        completed = translated_count + failed_count
        data: dict[str, Any] = {
            "translated": translated_count,
            "failed": failed_count,
            "engine": "model",
            "translation_profile": translation_profile,
            "source_chars_done": completed_source_chars,
            "source_chars_total": target_source_chars,
            "manifest_path": str(manifest_path.resolve()),
            "resumable": True,
        }
        if fallback:
            data["fallback"] = True
        if paragraph_index is not None:
            data["paragraph_index"] = paragraph_index
        context.log(
            "info",
            f"translation progress {completed}/{target_nonempty_goal}",
            data,
        )

    async def flush_pending() -> None:
        nonlocal pending_chars, completed_source_chars
        if not pending:
            return

        batch = list(pending)
        pending.clear()
        pending_chars = 0
        batch_source_chars = sum(int(item.get("source_chars") or 0) for item in batch)
        flat_texts: list[str] = []
        chunk_counts: list[int] = []
        for item in batch:
            paragraph_text = str(item.get("source_text") or "")
            chunks = _split_translation_chunks(paragraph_text, max_chars_per_chunk)
            chunk_counts.append(len(chunks))
            flat_texts.extend(chunks)

        start_index = int(batch[0].get("paragraph_index") or 0)
        end_index = int(batch[-1].get("paragraph_index") or 0)
        completed_before = translated_count + failed_count
        context.log(
            "info",
            f"translation batch started {completed_before}/{target_nonempty_goal}",
            {
                "paragraph_start": start_index,
                "paragraph_end": end_index,
                "batch_paragraphs": len(batch),
                "batch_chars": batch_source_chars,
                "chunk_count": len(flat_texts),
                "timeout": batch_timeout,
                "source_chars_done": completed_source_chars,
                "source_chars_total": target_source_chars,
                "engine": "model",
                "translation_profile": translation_profile,
                "manifest_path": str(manifest_path.resolve()),
            },
        )

        try:
            translated_chunks = await asyncio.wait_for(
                _translate_text_batch_with_model(
                    settings=settings,
                    model=model,
                    texts=flat_texts,
                    source_language=source_language,
                    target_language=target_language,
                ),
                timeout=batch_timeout,
            )
        except Exception as batch_exc:
            if len(batch) == 1:
                append_failed_paragraph(batch[0], batch_exc)
                completed_source_chars += batch_source_chars
                _save_translate_docx_manifest(manifest_path, manifest)
            else:
                update_manifest_warning(f"batch fallback: {str(batch_exc)[:200]}")
                for item in batch:
                    paragraph_text = str(item.get("source_text") or "")
                    try:
                        chunks = _split_translation_chunks(paragraph_text, max_chars_per_chunk)
                        translated_parts = await asyncio.wait_for(
                            _translate_text_batch_with_model(
                                settings=settings,
                                model=model,
                                texts=chunks,
                                source_language=source_language,
                                target_language=target_language,
                            ),
                            timeout=min(batch_timeout, 120),
                        )
                        translated = "".join(part or "" for part in translated_parts).strip()
                        if not translated:
                            raise RuntimeError("model returned empty text")
                        append_translated_paragraph(item, translated)
                    except Exception as exc:
                        append_failed_paragraph(item, exc)
                    completed_source_chars += len(paragraph_text)
                    _save_translate_docx_manifest(manifest_path, manifest)
                    completed = translated_count + failed_count
                    if completed and (completed % save_every == 0 or completed >= target_nonempty_goal):
                        target_doc.save(str(output_path))
                    log_translation_progress(
                        fallback=True,
                        paragraph_index=int(item.get("paragraph_index") or 0),
                    )
                return
            log_translation_progress()
            return

        position = 0
        for item, chunk_count in zip(batch, chunk_counts):
            translated = "".join(translated_chunks[position:position + chunk_count]).strip()
            position += chunk_count
            if not translated:
                append_failed_paragraph(item, RuntimeError("model returned empty text"))
                continue
            append_translated_paragraph(item, translated)

        completed_source_chars += batch_source_chars
        _save_translate_docx_manifest(manifest_path, manifest)
        completed = translated_count + failed_count
        if completed and (completed % save_every == 0 or completed >= target_nonempty_goal):
            target_doc.save(str(output_path))
        log_translation_progress()

    for item in manifest_items[processed_nonempty:]:
        if max_seconds and time.monotonic() - started_at > max_seconds:
            stopped_reason = f"max_seconds_exceeded:{max_seconds}"
            break

        if max_paragraphs and processed_nonempty >= max_paragraphs:
            stopped_reason = f"max_paragraphs_reached:{max_paragraphs}"
            break

        processed_nonempty += 1
        pending.append(item)
        pending_chars += int(item.get("source_chars") or 0)
        if len(pending) >= max_paragraphs_per_batch or pending_chars >= max_chars_per_batch:
            await flush_pending()

    if pending:
        await flush_pending()

    if not stopped_reason and max_paragraphs and target_nonempty_goal < source_nonempty:
        stopped_reason = f"max_paragraphs_reached:{max_paragraphs}"

    target_doc.save(str(output_path))

    complete = (
        translated_count == source_nonempty
        and failed_count == 0
        and not stopped_reason
    )
    partial_resumable = (
        not complete
        and stopped_reason.startswith("max_seconds_exceeded")
        and (translated_count + failed_count) < target_nonempty_goal
        and (translated_count + failed_count) > 0
    )
    status = "success" if complete else ("partial_resumable" if partial_resumable else "partial")
    output: dict[str, Any] = {
        "path": str(output_path.resolve()),
        "source_path": str(source_path.resolve()),
        "manifest_path": str(manifest_path.resolve()),
        "engine": "model",
        "model": model,
        "translation_profile": translation_profile,
        "resumed_from_checkpoint": resumed_from_checkpoint,
        "partial_resumable": partial_resumable,
        "resume": True,
        "resume_input": {
            "path": str(source_path.resolve()),
            "output_path": str(output_path.resolve()),
            "manifest_path": str(manifest_path.resolve()),
            "source_language": source_language,
            "target_language": target_language,
            "resume": True,
        },
        "source_language": source_language,
        "target_language": target_language,
        "source_paragraph_count": len(paragraphs),
        "source_nonempty_paragraph_count": source_nonempty,
        "target_nonempty_goal": target_nonempty_goal,
        "processed_paragraph_count": processed_nonempty,
        "translated_paragraph_count": translated_count,
        "failed_paragraph_count": failed_count,
        "source_chars_total": target_source_chars,
        "source_chars_done": completed_source_chars,
        "max_seconds": max_seconds,
        "batch_timeout": batch_timeout,
        "max_paragraphs_per_batch": max_paragraphs_per_batch,
        "max_chars_per_batch": max_chars_per_batch,
        "max_chars_per_chunk": max_chars_per_chunk,
        "translated_chars": translated_chars,
        "complete": complete,
        "stopped_reason": stopped_reason,
        "file_size": output_path.stat().st_size if output_path.exists() else 0,
        "warnings": warnings[:20],
        "status": status,
    }
    return output


def _translate_docx_google_sync(input_data: dict[str, Any], context: Any) -> dict[str, Any]:
    source_path = context.path_guard.resolve(
        input_data.get("path") or input_data.get("source_path")
    )
    if source_path.suffix.lower() != ".docx":
        raise ValueError("仅支持 .docx 格式")

    output_path = _resolve_translate_docx_output_path(input_data, context, source_path)
    source_language = str(input_data.get("source_language") or input_data.get("source") or "auto")
    target_language = str(input_data.get("target_language") or input_data.get("target") or "zh-CN")
    max_chars_per_chunk = max(500, min(int(input_data.get("max_chars_per_chunk", 3000)), 4500))
    max_paragraphs = max(0, int(input_data.get("max_paragraphs", 0) or 0))
    max_seconds = max(30, min(int(input_data.get("max_seconds", 600) or 600), 3600))
    save_every = max(1, min(int(input_data.get("save_every", 20) or 20), 100))

    try:
        from docx import Document
        from docx.oxml.ns import qn
        from docx.shared import Pt
    except ImportError as exc:
        raise RuntimeError("python-docx is required for document.translate_docx") from exc
    try:
        from deep_translator import GoogleTranslator
    except ImportError as exc:
        raise RuntimeError("deep-translator is required for document.translate_docx") from exc

    source_doc = Document(str(source_path))
    paragraphs = list(source_doc.paragraphs)
    source_nonempty = sum(1 for paragraph in paragraphs if paragraph.text.strip())
    target_nonempty_goal = min(source_nonempty, max_paragraphs) if max_paragraphs else source_nonempty

    target_doc = Document()
    normal = target_doc.styles["Normal"]
    normal.font.name = "Microsoft YaHei"
    normal._element.rPr.rFonts.set(qn("w:eastAsia"), "Microsoft YaHei")
    normal.font.size = Pt(10.5)

    translator = GoogleTranslator(source=source_language, target=target_language)
    started_at = time.monotonic()
    processed_nonempty = 0
    translated_count = 0
    failed_count = 0
    translated_chars = 0
    warnings: list[str] = []
    stopped_reason = ""

    _backup_output(output_path, context)

    for index, paragraph in enumerate(paragraphs, 1):
        if max_seconds and time.monotonic() - started_at > max_seconds:
            stopped_reason = f"max_seconds_exceeded:{max_seconds}"
            break

        text = paragraph.text.strip()
        if not text:
            target_doc.add_paragraph()
            continue

        if max_paragraphs and processed_nonempty >= max_paragraphs:
            stopped_reason = f"max_paragraphs_reached:{max_paragraphs}"
            break

        processed_nonempty += 1
        try:
            translated_parts = [
                translator.translate(chunk)
                for chunk in _split_translation_chunks(text, max_chars_per_chunk)
            ]
            translated = "".join(part or "" for part in translated_parts).strip()
            if not translated:
                raise RuntimeError("translator returned empty text")
            target_doc.add_paragraph(translated)
            translated_count += 1
            translated_chars += len(translated)
        except Exception as exc:
            failed_count += 1
            message = f"[翻译失败，保留原文] {text}"
            target_doc.add_paragraph(message)
            warnings.append(f"paragraph {index}: {str(exc)[:200]}")

        if processed_nonempty % save_every == 0:
            target_doc.save(str(output_path))
            context.log(
                "info",
                f"translation progress {processed_nonempty}/{target_nonempty_goal}",
                {"translated": translated_count, "failed": failed_count, "engine": "google"},
            )

    target_doc.save(str(output_path))

    complete = (
        translated_count == source_nonempty
        and failed_count == 0
        and not stopped_reason
    )
    output: dict[str, Any] = {
        "path": str(output_path.resolve()),
        "source_path": str(source_path.resolve()),
        "engine": "google",
        "source_language": source_language,
        "target_language": target_language,
        "source_paragraph_count": len(paragraphs),
        "source_nonempty_paragraph_count": source_nonempty,
        "target_nonempty_goal": target_nonempty_goal,
        "processed_paragraph_count": processed_nonempty,
        "translated_paragraph_count": translated_count,
        "failed_paragraph_count": failed_count,
        "translated_chars": translated_chars,
        "complete": complete,
        "stopped_reason": stopped_reason,
        "file_size": output_path.stat().st_size if output_path.exists() else 0,
        "warnings": warnings[:20],
    }
    if not complete:
        output["error"] = True
        output["status"] = "partial"
    else:
        output["status"] = "success"
    return output


async def generate_docx_from_outline(input_data: dict[str, Any], context: Any) -> dict[str, Any]:
    title = input_data.get("title", "文档")
    path = _resolve_output_path(input_data, context, title, ".docx")
    outline = input_data.get("outline", [])

    try:
        from docx import Document
        from docx.oxml.ns import qn
        from docx.shared import Pt
    except ImportError as exc:
        raise RuntimeError("python-docx is required for document.generate_docx_from_outline") from exc

    doc = Document()

    normal = doc.styles["Normal"]
    normal.font.name = "Microsoft YaHei"
    normal._element.rPr.rFonts.set(qn("w:eastAsia"), "Microsoft YaHei")
    normal.font.size = Pt(10.5)

    for style_name in ("Heading 1", "Heading 2", "Heading 3"):
        style = doc.styles[style_name]
        style.font.name = "Microsoft YaHei"
        style._element.rPr.rFonts.set(qn("w:eastAsia"), "Microsoft YaHei")

    doc.add_heading(title, level=1)

    for item in outline:
        level = min(int(item.get("level", 1)), 3)
        text = item.get("text", "")
        if text:
            doc.add_heading(text, level=level)

    _backup_output(path, context)
    doc.save(str(path))

    return {
        "path": str(path.resolve()),
        "title": title,
        "outline_count": len(outline),
        "content_chars": sum(len(str(item.get("text") or "")) for item in outline),
        "paragraph_count": len(doc.paragraphs),
        "nonempty_paragraph_count": sum(1 for paragraph in doc.paragraphs if paragraph.text.strip()),
        "file_size": path.stat().st_size if path.exists() else 0,
    }


async def export_pdf(input_data: dict[str, Any], context: Any) -> dict[str, Any]:
    title = input_data.get("title", "AI生成文档")
    path = _resolve_output_path(input_data, context, title, ".pdf")
    content = input_data.get("content", "")

    try:
        from markdown import markdown
    except ImportError as exc:
        raise RuntimeError("markdown library is required for document.export_pdf") from exc

    html_content = markdown(content)

    try:
        from playwright.async_api import async_playwright
    except ImportError as exc:
        raise RuntimeError(
            "playwright is required for document.export_pdf. "
            "Install: pip install playwright && playwright install chromium"
        ) from exc

    full_html = f"""<!DOCTYPE html>
<html lang="zh-CN">
<head>
    <meta charset="utf-8">
    <title>{title}</title>
    <style>
        body {{ font-family: "Microsoft YaHei", "PingFang SC", Arial, sans-serif; color: #111827; line-height: 1.72; font-size: 14px; max-width: 800px; margin: 0 auto; padding: 20px; }}
        h1, h2, h3, h4 {{ color: #0f172a; line-height: 1.35; margin: 20px 0 10px; }}
        h1 {{ font-size: 22px; border-bottom: 2px solid #e5e7eb; padding-bottom: 8px; }}
        h2 {{ font-size: 18px; }}
        h3 {{ font-size: 16px; }}
        pre {{ background: #f3f4f6; padding: 12px; border-radius: 6px; overflow-x: auto; font-size: 13px; }}
        code {{ font-family: Consolas, "Courier New", monospace; background: #f3f4f6; padding: 2px 4px; border-radius: 3px; font-size: 13px; }}
        pre code {{ background: none; padding: 0; }}
        blockquote {{ border-left: 4px solid #d1d5db; margin: 12px 0; padding: 8px 16px; color: #6b7280; }}
        table {{ border-collapse: collapse; width: 100%; margin: 12px 0; }}
        th, td {{ border: 1px solid #d1d5db; padding: 8px 12px; text-align: left; }}
        th {{ background: #f9fafb; font-weight: 600; }}
        img {{ max-width: 100%; }}
    </style>
</head>
<body>
    <h1>{title}</h1>
    {html_content}
</body>
</html>"""

    async with async_playwright() as playwright:
        browser = await playwright.chromium.launch(headless=True)
        try:
            page = await browser.new_page(viewport={"width": 1240, "height": 1754})
            await page.set_content(full_html, wait_until="load")
            _backup_output(path, context)
            await page.pdf(
                path=str(path),
                format="A4",
                print_background=True,
                margin={"top": "18mm", "right": "18mm", "bottom": "18mm", "left": "18mm"},
            )
        finally:
            await browser.close()

    return {
        "path": str(path.resolve()),
        "title": title,
    }


async def generate_ppt(input_data: dict[str, Any], context: Any) -> dict[str, Any]:
    title = input_data.get("title", "演示文稿")
    path = _resolve_output_path(input_data, context, title, ".pptx")

    slides = input_data.get("slides", [])
    outline = input_data.get("outline", [])
    content = input_data.get("content", "")

    # 兼容多种输入格式：优先slides，其次outline，最后纯内容
    if not slides and outline:
        # 从大纲自动生成幻灯片
        slides = []
        # 封面页
        slides.append({"title": title, "content": "AI 生成演示文稿"})
        # 内容页
        current_slide = {"title": "", "content": []}
        for item in outline:
            level = item.get("level", 1)
            text = item.get("text", "").strip()
            if not text:
                continue
            if level == 1:
                if current_slide["title"]:
                    slides.append({
                        "title": current_slide["title"],
                        "content": "\n- " + "\n- ".join(current_slide["content"])
                    })
                current_slide = {"title": text, "content": []}
            else:
                current_slide["content"].append(text)
        if current_slide["title"]:
            slides.append({
                "title": current_slide["title"],
                "content": "\n- " + "\n- ".join(current_slide["content"])
            })
    elif not slides and content:
        # 从纯文本自动生成幻灯片
        lines = [line.strip() for line in content.split("\n") if line.strip()]
        slides = []
        # 封面页
        slides.append({"title": title, "content": "AI 生成演示文稿"})
        # 内容页
        current_title = ""
        current_content = []
        for line in lines:
            if line.startswith("# "):
                if current_title:
                    slides.append({
                        "title": current_title,
                        "content": "\n".join(current_content)
                    })
                current_title = line[2:].strip()
                current_content = []
            elif line.startswith("## "):
                if current_title:
                    slides.append({
                        "title": current_title,
                        "content": "\n".join(current_content)
                    })
                current_title = line[3:].strip()
                current_content = []
            else:
                current_content.append(line)
        if current_title:
            slides.append({
                "title": current_title,
                "content": "\n".join(current_content)
            })

    if not slides:
        raise ValueError("slides/outline/content is required: provide at least one content source")

    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
        from pptx.dml.color import RGBColor
    except ImportError as exc:
        raise RuntimeError("python-pptx is required for document.generate_ppt") from exc

    prs = Presentation()

    # 设置默认字体为微软雅黑，支持中文
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text_frame"):
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        run.font.name = "Microsoft YaHei"
                        run.font.element.rPr.rFonts.set(prs.nsdecls['w'], "Microsoft YaHei")

    for i, slide_data in enumerate(slides):
        if i == 0:
            # 封面页用布局0
            slide_layout = prs.slide_layouts[0]
        else:
            # 内容页用布局1，带标题和内容
            slide_layout = prs.slide_layouts[1]

        slide = prs.slides.add_slide(slide_layout)

        title_text = slide_data.get("title", "")
        if title_text and slide.shapes.title:
            title_shape = slide.shapes.title
            title_shape.text = title_text
            # 标题字体设置
            for paragraph in title_shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.name = "Microsoft YaHei"
                    run.font.element.rPr.rFonts.set(prs.nsdecls['w'], "Microsoft YaHei")
                    run.font.size = Pt(24 if i == 0 else 20)
                    if i == 0:
                        run.font.bold = True

        content_text = slide_data.get("content", "")
        if content_text and slide.shapes.placeholders:
            # 找到内容占位符（idx=1）
            content_placeholder = None
            for placeholder in slide.shapes.placeholders:
                if placeholder.placeholder_format.idx == 1:
                    content_placeholder = placeholder
                    break

            if content_placeholder:
                content_placeholder.text = content_text
                # 内容字体设置
                for paragraph in content_placeholder.text_frame.paragraphs:
                    for run in paragraph.runs:
                        run.font.name = "Microsoft YaHei"
                        run.font.element.rPr.rFonts.set(prs.nsdecls['w'], "Microsoft YaHei")
                        run.font.size = Pt(14)

    _backup_output(path, context)
    prs.save(str(path))

    return {
        "path": str(path.resolve()),
        "title": title,
        "slide_count": len(slides),
        "input_type": "slides" if input_data.get("slides") else "outline" if input_data.get("outline") else "content"
    }


async def merge_pdfs(input_data: dict[str, Any], context: Any) -> dict[str, Any]:
    output_path = context.path_guard.resolve(input_data.get("output_path"))
    input_paths = input_data.get("input_paths", [])

    try:
        from pypdf import PdfMerger
    except ImportError as exc:
        raise RuntimeError("pypdf is required for document.merge_pdfs") from exc

    merger = PdfMerger()

    for path_str in input_paths:
        path = context.path_guard.resolve(path_str)
        merger.append(str(path))

    _backup_output(output_path, context)
    try:
        merger.write(str(output_path))
    finally:
        merger.close()

    return {
        "output_path": str(output_path.resolve()),
        "merged_count": len(input_paths),
    }


async def split_pdf(input_data: dict[str, Any], context: Any) -> dict[str, Any]:
    path = context.path_guard.resolve(input_data.get("path"))
    output_dir = context.path_guard.resolve(input_data.get("output_dir", path.parent))

    try:
        from pypdf import PdfReader, PdfWriter
    except ImportError as exc:
        raise RuntimeError("pypdf is required for document.split_pdf") from exc

    reader = PdfReader(str(path))
    output_dir.mkdir(exist_ok=True)

    split_files = []

    for i, page in enumerate(reader.pages):
        writer = PdfWriter()
        writer.add_page(page)

        output_path = output_dir / f"{path.stem}_page_{i+1}.pdf"
        _backup_output(output_path, context)
        with open(output_path, "wb") as f:
            writer.write(f)

        split_files.append(str(output_path))

    return {
        "source_path": str(path),
        "output_dir": str(output_dir),
        "split_files": split_files,
    }


async def create_bookmark_outline(input_data: dict[str, Any], context: Any) -> dict[str, Any]:
    path = context.path_guard.resolve(input_data.get("path"))
    bookmarks = input_data.get("bookmarks", [])

    try:
        from pypdf import PdfReader, PdfWriter
    except ImportError as exc:
        raise RuntimeError("pypdf is required for document.create_bookmark_outline") from exc

    reader = PdfReader(str(path))
    writer = PdfWriter()

    for page in reader.pages:
        writer.add_page(page)

    def add_bookmarks(bookmarks_list: list[dict[str, Any]], parent: Any = None) -> None:
        for item in bookmarks_list:
            title = item.get("title", "")
            page_num = max(0, int(item.get("page", 1)) - 1)
            if page_num < len(reader.pages):
                dest = writer.add_outline_item(title, page_num, parent)
                children = item.get("children", [])
                if children:
                    add_bookmarks(children, dest)

    add_bookmarks(bookmarks)

    output_path = path.parent / f"{path.stem}_with_bookmarks.pdf"
    _backup_output(output_path, context)
    with open(output_path, "wb") as f:
        writer.write(f)

    return {
        "output_path": str(output_path.resolve()),
        "bookmark_count": len(bookmarks),
    }


def register_document_tools(registry: ToolRegistry) -> None:
    registry.register(
        ToolSpec(
            id="document.extract_docx_outline",
            name="智能提取 Word 文档",
            description="智能提取 Word 文档内容（大纲 + 全文 + 表格）。支持 .docx 和 .doc 格式，自动处理损坏/非标准文件。",
            input_schema={
                "type": "object",
                "properties": {
                    "path": {"type": "string", "description": "Word 文档路径（.docx 或 .doc）"},
                    "extract_text": {"type": "boolean", "default": True, "description": "是否提取全文内容（包含段落和表格）"},
                },
                "required": ["path"],
            },
            optional_dependencies=["docx"],
        ),
        extract_docx_outline,
    )
    registry.register(
        ToolSpec(
            id="document.extract_pdf_text_preview",
            name="智能提取 PDF 文本",
            description="智能提取 PDF 文本内容（自动乱码检测 + 多引擎兜底 + AI OCR 回退）。支持扫描版和加密字体 PDF。",
            input_schema={
                "type": "object",
                "properties": {
                    "path": {"type": "string", "description": "PDF 文件路径"},
                    "max_pages": {"type": "integer", "default": 50, "description": "最大解析页数，0表示不限"},
                },
                "required": ["path"],
            },
            optional_dependencies=["pypdf"],
        ),
        extract_pdf_text_preview,
    )
    registry.register(
        ToolSpec(
            id="document.extract_pdf_to_docx",
            name="PDF 文本转存 Word",
            description="提取 PDF 文本并直接保存为 .docx，适合 PDF 转 Word、PDF 文本转存 Word 等任务。",
            input_schema={
                "type": "object",
                "properties": {
                    "path": {"type": "string", "description": "PDF 文件路径"},
                    "output_path": {"type": "string", "description": "输出 Word 路径（可选，默认在 PDF 同目录生成 *_extracted.docx）"},
                    "mode": {
                        "type": "string",
                        "default": "text_only",
                        "enum": ["text_only", "text_with_images"],
                        "description": "PDF 转 Word 模式：text_only 只提取文字；text_with_images 按页面块顺序近似保留文字和图片",
                    },
                    "max_pages": {"type": "integer", "default": 0, "description": "最大解析页数，0表示不限"},
                    "title": {"type": "string", "description": "Word 文档标题（可选）"},
                },
                "required": ["path"],
            },
            requires_confirmation=True,
            optional_dependencies=["pypdf", "docx", "fitz", "PIL"],
            capability="document.pdf_to_docx",
            artifacts=["docx"],
            long_running=True,
            retry_safe=True,
        ),
        extract_pdf_to_docx,
    )
    registry.register(
        ToolSpec(
            id="document.export_markdown",
            name="导出 Markdown 文档",
            description="将文本内容保存为 .md 文件，会自动添加标题。",
            input_schema={
                "type": "object",
                "properties": {
                    "path": {"type": "string"},
                    "content": {"type": "string"},
                    "title": {"type": "string", "default": "AI生成文档"},
                },
                "required": ["path", "content"],
            },
            requires_confirmation=True,
            capability="document.export_markdown",
            artifacts=["markdown"],
            idempotent=True,
        ),
        export_markdown,
    )
    registry.register(
        ToolSpec(
            id="document.export_docx",
            name="导出 Word 文档",
            description="将 Markdown 内容导出为 .docx 文档，支持基本的标题、段落和代码块。",
            input_schema={
                "type": "object",
                "properties": {
                    "path": {"type": "string"},
                    "content": {"type": "string"},
                    "title": {"type": "string", "default": "AI生成文档"},
                },
                "required": ["path", "content"],
            },
            requires_confirmation=True,
            optional_dependencies=["docx"],
            capability="document.export_docx",
            artifacts=["docx"],
            idempotent=True,
        ),
        export_docx,
    )
    registry.register(
        ToolSpec(
            id="document.create_draft",
            name="Create document draft",
            description=(
                "Create a persistent document draft state. Use this for long reports, papers, "
                "book manuscripts, long translations, or any task that needs progressive writing "
                "before a final export. When expanding an existing Word document, pass its .docx "
                "path as source_path so the draft starts with the existing content instead of empty."
            ),
            input_schema={
                "type": "object",
                "properties": {
                    "title": {"type": "string"},
                    "source_path": {
                        "type": "string",
                        "description": "Optional existing .docx file to import before appending new content.",
                    },
                    "sections": {
                        "type": "array",
                        "items": {
                            "type": "object",
                            "properties": {
                                "section_id": {"type": "string"},
                                "title": {"type": "string"},
                                "level": {"type": "integer", "default": 1},
                                "metadata": {"type": "object"},
                            },
                        },
                    },
                    "metadata": {"type": "object"},
                },
            },
            capability="document.draft",
            artifacts=["draft"],
            retry_safe=True,
        ),
        create_draft,
    )
    registry.register(
        ToolSpec(
            id="document.append_draft_section",
            name="Append document draft content",
            description=(
                "Append content to a document draft section, creating the section when needed. "
                "This is the preferred write path for long-form content because the draft can be "
                "inspected, resumed, and exported later. Keep each call complete and bounded; "
                "use several calls when one complete block would exceed the model output budget."
            ),
            input_schema={
                "type": "object",
                "properties": {
                    "draft_id": {"type": "string"},
                    "section_id": {"type": "string"},
                    "title": {"type": "string"},
                    "level": {"type": "integer", "default": 1},
                    "content": {
                        "type": "string",
                        "description": "One complete content block.",
                    },
                    "citation_ids": {"type": "array", "items": {"type": "string"}},
                    "metadata": {"type": "object"},
                },
                "required": ["draft_id", "content"],
            },
            capability="document.draft",
            artifacts=["draft"],
            retry_safe=True,
        ),
        append_draft_section,
    )
    registry.register(
        ToolSpec(
            id="document.add_draft_citation",
            name="Add document draft citation",
            description="Add a citation/source record to a document draft for papers, reports, or book notes.",
            input_schema={
                "type": "object",
                "properties": {
                    "draft_id": {"type": "string"},
                    "citation": {
                        "type": "object",
                        "properties": {
                            "citation_id": {"type": "string"},
                            "title": {"type": "string"},
                            "source_type": {"type": "string"},
                            "url": {"type": "string"},
                            "doi": {"type": "string"},
                            "author": {"type": "string"},
                            "year": {"type": "string"},
                            "note": {"type": "string"},
                        },
                    },
                },
                "required": ["draft_id", "citation"],
            },
            capability="document.draft",
            artifacts=["draft"],
            retry_safe=True,
        ),
        add_draft_citation,
    )
    registry.register(
        ToolSpec(
            id="document.inspect_draft",
            name="Inspect document draft",
            description="Inspect document draft structure, progress, citations, and missing references without exporting.",
            input_schema={
                "type": "object",
                "properties": {
                    "draft_id": {"type": "string"},
                },
                "required": ["draft_id"],
            },
            capability="document.draft",
            artifacts=["draft"],
        ),
        inspect_draft,
    )
    registry.register(
        ToolSpec(
            id="document.export_draft_docx",
            name="Export document draft to Word",
            description="Export a persistent document draft to .docx. The output path is protected by PathGuard and confirmation.",
            input_schema={
                "type": "object",
                "properties": {
                    "draft_id": {"type": "string"},
                    "path": {"type": "string"},
                    "output_path": {"type": "string"},
                    "include_citations": {"type": "boolean", "default": True},
                },
                "required": ["draft_id"],
            },
            requires_confirmation=True,
            optional_dependencies=["docx"],
            capability="document.draft_export",
            artifacts=["docx", "draft"],
            retry_safe=True,
            idempotent=True,
        ),
        export_draft_docx,
    )
    registry.register(
        ToolSpec(
            id="document.translate_docx",
            name="翻译 Word 文档",
            description=(
                "逐段翻译 .docx 文档并保存为新的 .docx。适合全文翻译、生成中文版。"
                "会返回翻译段落数、失败段落数和是否完整完成；部分完成不会被当作完整成功。"
            ),
            input_schema={
                "type": "object",
                "properties": {
                    "path": {"type": "string", "description": "源 Word 文档路径（.docx）"},
                    "output_path": {"type": "string", "description": "输出 Word 路径（可选）"},
                    "engine": {
                        "type": "string",
                        "default": "model",
                        "enum": ["model"],
                        "description": "翻译引擎，默认且推荐使用当前配置的大模型",
                    },
                    "model": {"type": "string", "description": "engine=model 时可指定模型 ID；为空则使用默认模型"},
                    "source_language": {"type": "string", "default": "auto", "description": "源语言，默认 auto"},
                    "target_language": {"type": "string", "default": "zh-CN", "description": "目标语言，默认 zh-CN"},
                    "manifest_path": {"type": "string", "description": "翻译 checkpoint JSON 路径；为空则默认生成在输出 Word 同目录"},
                    "resume": {"type": "boolean", "default": True, "description": "是否从匹配的 checkpoint 继续翻译"},
                    "reset_checkpoint": {"type": "boolean", "default": False, "description": "忽略已有 checkpoint 并重新开始"},
                    "translation_profile": {
                        "type": "string",
                        "default": "balanced",
                        "enum": ["safe", "balanced", "fast"],
                        "description": "翻译分批档位：balanced 默认保持较快进度回显；fast 适合大上下文模型；safe 适合容易超时的模型",
                    },
                    "max_paragraphs": {"type": "integer", "default": 0, "description": "最多翻译的非空段落数，0 表示全部"},
                    "max_seconds": {"type": "integer", "default": 1800, "description": "最长运行秒数，全文翻译会自动使用足够的最低预算"},
                    "max_paragraphs_per_batch": {"type": "integer", "default": 8, "description": "每批最多翻译段落数"},
                    "max_chars_per_batch": {"type": "integer", "default": 12000, "description": "每批最多翻译字符数"},
                    "max_chars_per_chunk": {"type": "integer", "default": 4000, "description": "单次翻译最大字符数"},
                    "batch_timeout": {"type": "integer", "default": 240, "description": "单批模型翻译最长等待秒数"},
                },
                "required": ["path"],
            },
            requires_confirmation=True,
            optional_dependencies=["docx"],
            capability="document.translate_docx",
            artifacts=["docx", "checkpoint"],
            long_running=True,
            retry_safe=True,
        ),
        translate_docx,
    )
    registry.register(
        ToolSpec(
            id="document.generate_docx_from_outline",
            name="从大纲生成 Word",
            description="根据大纲标题列表快速生成 .docx 文档框架。",
            input_schema={
                "type": "object",
                "properties": {
                    "path": {"type": "string"},
                    "title": {"type": "string", "default": "文档"},
                    "output_path": {"type": "string"},
                    "outline": {
                        "type": "array",
                        "items": {
                            "type": "object",
                            "properties": {
                                "level": {"type": "integer", "default": 1},
                                "text": {"type": "string"},
                            },
                        },
                    },
                },
                "required": ["outline"],
            },
            requires_confirmation=True,
            optional_dependencies=["docx"],
            capability="document.generate_docx",
            artifacts=["docx"],
            idempotent=True,
        ),
        generate_docx_from_outline,
    )
    registry.register(
        ToolSpec(
            id="document.export_pdf",
            name="导出 PDF 文档",
            description="将 Markdown 内容导出为 .pdf 文档。需要安装 playwright。",
            input_schema={
                "type": "object",
                "properties": {
                    "path": {"type": "string"},
                    "content": {"type": "string"},
                    "title": {"type": "string", "default": "AI生成文档"},
                },
                "required": ["path", "content"],
            },
            requires_confirmation=True,
            optional_dependencies=["markdown", "playwright"],
            capability="document.export_pdf",
            artifacts=["pdf"],
            long_running=True,
            idempotent=True,
            readiness_probe=playwright_chromium_readiness,
        ),
        export_pdf,
    )
    registry.register(
        ToolSpec(
            id="document.generate_ppt",
            name="生成 PowerPoint",
            description="根据幻灯片内容列表快速生成 .pptx 演示文稿。每张幻灯片需要 title（标题）和 content（正文，多条用换行分隔）。第一张使用封面版式，其余使用标题+内容版式。",
            optional_dependencies=["pptx"],
            input_schema={
                "type": "object",
                "properties": {
                    "path": {"type": "string", "description": "输出的 .pptx 文件路径"},
                    "title": {"type": "string", "default": "演示文稿", "description": "演示文稿总标题"},
                    "slides": {
                        "type": "array",
                        "description": "幻灯片列表，每项为一张幻灯片",
                        "items": {
                            "type": "object",
                            "properties": {
                                "title": {"type": "string", "description": "幻灯片标题"},
                                "content": {"type": "string", "description": "幻灯片正文内容，多条要点用换行符分隔"},
                            },
                            "required": ["title"],
                        },
                    },
                },
                "required": ["path", "slides"],
            },
            requires_confirmation=True,
            capability="document.generate_ppt",
            artifacts=["pptx"],
            idempotent=True,
        ),
        generate_ppt,
    )
    registry.register(
        ToolSpec(
            id="document.merge_pdfs",
            name="合并 PDF 文件",
            description="将多个 PDF 文件合并为一个。",
            input_schema={
                "type": "object",
                "properties": {
                    "output_path": {"type": "string"},
                    "input_paths": {"type": "array", "items": {"type": "string"}},
                },
                "required": ["output_path", "input_paths"],
            },
            requires_confirmation=True,
            optional_dependencies=["pypdf"],
            capability="document.merge_pdf",
            artifacts=["pdf"],
            idempotent=True,
        ),
        merge_pdfs,
    )
    registry.register(
        ToolSpec(
            id="document.split_pdf",
            name="拆分 PDF 文件",
            description="将 PDF 按页拆分为多个文件。",
            input_schema={
                "type": "object",
                "properties": {
                    "path": {"type": "string"},
                    "output_dir": {"type": "string"},
                },
                "required": ["path"],
            },
            requires_confirmation=True,
            optional_dependencies=["pypdf"],
        ),
        split_pdf,
    )
    registry.register(
        ToolSpec(
            id="document.create_bookmark_outline",
            name="创建 PDF 书签",
            description="为 PDF 文件添加书签导航。",
            input_schema={
                "type": "object",
                "properties": {
                    "path": {"type": "string"},
                    "bookmarks": {
                        "type": "array",
                        "items": {
                            "type": "object",
                            "properties": {
                                "title": {"type": "string"},
                                "page": {"type": "integer"},
                                "children": {"type": "array", "items": {"type": "object"}},
                            },
                        },
                    },
                },
                "required": ["path", "bookmarks"],
            },
            requires_confirmation=True,
            optional_dependencies=["pypdf"],
        ),
        create_bookmark_outline,
    )
